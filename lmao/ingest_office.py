import argparse
import os
import requests
import csv
import json
from msal import ConfidentialClientApplication
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
from io import BytesIO
from typing import List, Dict, Any


class OfficeIngester:
    def __init__(
        self,
        client_id: str,
        client_secret: str,
        tenant_id: str,
        site_id: str,
        user_email: str,
        download_dir: str,
        batch_size: int = 10,
    ) -> None:
        """
        Initialize the OfficeIngester with necessary credentials and configurations.

        :param client_id: The client ID of the Azure AD app.
        :param client_secret: The client secret of the Azure AD app.
        :param tenant_id: The tenant ID of the Azure AD app.
        :param site_id: The ID of the SharePoint site.
        :param user_email: The email of the user whose contributions to search for.
        :param download_dir: The directory to download and save the documents.
        :param batch_size: The number of documents to process in each batch.
        """
        self.client_id = client_id
        self.client_secret = client_secret
        self.tenant_id = tenant_id
        self.site_id = site_id
        self.user_email = user_email
        self.download_dir = download_dir
        self.batch_size = batch_size
        self.access_token = self.get_access_token()

    def get_access_token(self) -> str:
        """
        Obtain an access token from Azure AD.

        :return: Access token as a string.
        :raises Exception: If the access token cannot be obtained.
        """
        authority = f"https://login.microsoftonline.com/{self.tenant_id}"
        app = ConfidentialClientApplication(
            self.client_id, authority=authority, client_credential=self.client_secret
        )

        result = app.acquire_token_for_client(
            scopes=["https://graph.microsoft.com/.default"]
        )

        if "access_token" in result:
            return result["access_token"]
        else:
            raise Exception("Could not obtain access token")

    def list_sharepoint_documents(self) -> Dict[str, Any]:
        """
        List documents in the SharePoint site.

        :return: A dictionary containing the list of documents.
        """
        headers = {"Authorization": f"Bearer {self.access_token}"}
        drive_items_url = (
            f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/drive/root/children"
        )

        response = requests.get(drive_items_url, headers=headers)
        response.raise_for_status()
        return response.json()

    def get_document_versions(self, item_id: str) -> Dict[str, Any]:
        """
        Get versions of a specific document.

        :param item_id: The ID of the document.
        :return: A dictionary containing the versions of the document.
        """
        headers = {"Authorization": f"Bearer {self.access_token}"}
        versions_url = f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/drive/items/{item_id}/versions"

        response = requests.get(versions_url, headers=headers)
        response.raise_for_status()
        return response.json()

    def filter_documents_by_user(
        self, drive_items: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """
        Filter documents by the specified user.

        :param drive_items: A dictionary containing the list of documents.
        :return: A list of filtered documents.
        """
        filtered_items = []
        for item in drive_items["value"]:
            versions = self.get_document_versions(item["id"])
            for version in versions["value"]:
                if (
                    "lastModifiedBy" in version
                    and version["lastModifiedBy"]["user"]["email"] == self.user_email
                ):
                    filtered_items.append(item)
                    break
        return filtered_items

    def download_document(self, file_id: str) -> bytes:
        """
        Download a document from SharePoint.

        :param file_id: The ID of the document to download.
        :return: The content of the document as bytes.
        """
        headers = {"Authorization": f"Bearer {self.access_token}"}
        download_url = f"https://graph.microsoft.com/v1.0/sites/{self.site_id}/drive/items/{file_id}/content"

        response = requests.get(download_url, headers=headers)
        response.raise_for_status()
        return response.content

    def extract_text_from_docx(self, file_path: str) -> str:
        """
        Extract text from a DOCX file.

        :param file_path: The path to the DOCX file.
        :return: Extracted text as a string.
        """
        doc = Document(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return "\n".join(full_text)

    def extract_images_from_docx(self, file_path: str) -> List[str]:
        """
        Extract images from a DOCX file.

        :param file_path: The path to the DOCX file.
        :return: A list of paths to the extracted images.
        """
        doc = Document(file_path)
        images = []

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img = rel.target_part.blob
                img_name = os.path.basename(rel.target_ref)
                img_path = os.path.join(self.download_dir, img_name)

                with open(img_path, "wb") as img_file:
                    img_file.write(img)
                images.append(img_path)

        return images

    def extract_text_from_pptx(self, file_path: str) -> str:
        """
        Extract text from a PPTX file.

        :param file_path: The path to the PPTX file.
        :return: Extracted text as a string.
        """
        prs = Presentation(file_path)
        full_text = [
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ]
        return "\n".join(full_text)

    def extract_images_from_pptx(self, file_path: str) -> List[str]:
        """
        Extract images from a PPTX file.

        :param file_path: The path to the PPTX file.
        :return: A list of paths to the extracted images.
        """
        prs = Presentation(file_path)
        images = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    img = Image.open(BytesIO(image.blob))
                    img_name = f"{shape.shape_id}.png"
                    img_path = os.path.join(self.download_dir, img_name)
                    img.save(img_path)
                    images.append(img_path)

        return images

    def extract_table_from_xlsx(
        self, file_path: str
    ) -> List[Dict[str, List[List[Any]]]]:
        """
        Extract tables from an XLSX file.

        :param file_path: The path to the XLSX file.
        :return: A list of dictionaries containing table data.
        """
        wb = load_workbook(file_path)
        table_data = []

        for sheet in wb.sheetnames:
            ws = wb[sheet]
            sheet_data = [list(row) for row in ws.iter_rows(values_only=True)]
            table_data.append({sheet: sheet_data})

        return table_data

    def process_batch(self, documents: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Process a batch of documents.

        :param documents: A list of documents to process.
        :return: A list of dictionaries containing document metadata and content.
        """
        batch_results = []

        for item in documents:
            file_id = item["id"]
            file_name = item["name"]
            file_content = self.download_document(file_id)
            file_path = os.path.join(self.download_dir, file_name)

            with open(file_path, "wb") as file:
                file.write(file_content)
            print(f"Downloaded: {file_name}")

            # Extract text and images based on file type
            if file_name.endswith(".docx"):
                text = self.extract_text_from_docx(file_path)
                images = self.extract_images_from_docx(file_path)
                content = {"text": text, "images": images}
            elif file_name.endswith(".pptx"):
                text = self.extract_text_from_pptx(file_path)
                images = self.extract_images_from_pptx(file_path)
                content = {"text": text, "images": images}
            elif file_name.endswith(".xlsx"):
                table_data = self.extract_table_from_xlsx(file_path)
                content = {"table": table_data}
            else:
                content = {}

            document_metadata = {
                "name": file_name,
                "id": file_id,
                "content": content,
                "metadata": item,
            }

            batch_results.append(document_metadata)

        return batch_results

    def run(self) -> None:
        """
        Run the OfficeIngester to download and process documents from SharePoint.
        """
        drive_items = self.list_sharepoint_documents()
        documents = self.filter_documents_by_user(drive_items)

        if not os.path.exists(self.download_dir):
            os.makedirs(self.download_dir)

        for i in range(0, len(documents), self.batch_size):
            batch = documents[i : i + self.batch_size]
            batch_results = self.process_batch(batch)
            batch_file_path = os.path.join(
                self.download_dir, f"batch_{i // self.batch_size + 1}.json"
            )

            with open(batch_file_path, "w") as batch_file:
                json.dump(batch_results, batch_file, indent=4)
            print(
                f"Processed batch {i // self.batch_size + 1} and saved to {batch_file_path}"
            )


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Download and extract documents from a SharePoint site."
    )

    parser.add_argument(
        "client_id", type=str, help="The client ID of the Azure AD app."
    )
    parser.add_argument(
        "client_secret", type=str, help="The client secret of the Azure AD app."
    )
    parser.add_argument(
        "tenant_id", type=str, help="The tenant ID of the Azure AD app."
    )
    parser.add_argument("site_id", type=str, help="The ID of the SharePoint site.")
    parser.add_argument(
        "user_email",
        type=str,
        help="The email of the user whose contributions to search for.",
    )
    parser.add_argument(
        "download_dir",
        type=str,
        help="The directory to download and save the documents.",
    )
    parser.add_argument(
        "--batch_size",
        type=int,
        default=10,
        help="The number of documents to process in each batch.",
    )

    args = parser.parse_args()

    ingester = OfficeIngester(
        client_id=args.client_id,
        client_secret=args.client_secret,
        tenant_id=args.tenant_id,
        site_id=args.site_id,
        user_email=args.user_email,
        download_dir=args.download_dir,
        batch_size=args.batch_size,
    )

    ingester.run()
